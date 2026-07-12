"""PPT exporter — generates a PowerPoint presentation from pipeline outputs.

Converts generated SVGs to PNG (via cairosvg) and embeds them into slides
with metadata (prompt, intent, score, key NLP metrics).
"""

import json
import logging
import os
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

OUTPUTS_DIR = Path(__file__).parent.parent.parent / "outputs"

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _svg_to_png_bytes(svg_code: str) -> bytes | None:
    """Convert SVG to PNG bytes using cairosvg."""
    try:
        import cairosvg

        return cairosvg.svg2png(bytestring=svg_code.encode("utf-8"))
    except Exception as e:
        logger.warning(f"SVG to PNG conversion failed: {e}")
        return None


def _add_title_slide(prs: Presentation) -> None:
    """Add a title/intro slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Multi-Agent SVG Generation System"
    if slide.placeholders[1]:
        slide.placeholders[1].text = (
            "基于多智能体的自然语言到 SVG 信息图生成系统\n\n"
            "Phase 3: 知识增强 + 渲染验证 + 反馈闭环\n"
            "2026-07-10"
        )


def _add_sample_slide(
    prs: Presentation,
    sample_name: str,
    prompt: str,
    metadata: dict[str, Any],
    svg_code: str,
) -> None:
    """Add a slide for one sample with SVG image and info."""
    # Use blank layout
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # ── Title ─────────────────────────────────────────────
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(12), Inches(0.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Sample: {sample_name}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # ── Prompt ────────────────────────────────────────────
    prompt_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8), Inches(8), Inches(0.5)
    )
    tf = prompt_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Prompt: {prompt}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # ── Metrics ───────────────────────────────────────────
    ir_summary = metadata.get("content_ir_summary", {})
    metrics_text = (
        f"Intent: {ir_summary.get('intent', '?')} | "
        f"Chart: {ir_summary.get('chart_type', '?')} | "
        f"Confidence: {ir_summary.get('confidence', '?')}\n"
        f"Score: {metadata.get('final_score', '?')} | "
        f"Passed: {metadata.get('passed', False)} | "
        f"Rounds: {metadata.get('refinement_rounds', 0)} | "
        f"Duration: {metadata.get('total_duration_s', 0):.0f}s"
    )
    metrics_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(8), Inches(0.5)
    )
    tf = metrics_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = metrics_text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)

    # ── SVG Image ─────────────────────────────────────────
    png_bytes = _svg_to_png_bytes(svg_code)
    if png_bytes:
        import io

        img_stream = io.BytesIO(png_bytes)
        try:
            slide.shapes.add_picture(
                img_stream,
                Inches(0.3),
                Inches(1.9),
                Inches(12.5),
                Inches(5.3),
            )
        except Exception as e:
            logger.warning(f"Failed to add picture: {e}")
    else:
        # Fallback: show text note
        note_box = slide.shapes.add_textbox(
            Inches(2), Inches(3.5), Inches(9), Inches(1)
        )
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[SVG image — cairosvg rendering not available]"
        p.font.size = Pt(16)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER


def export_ppt(output_path: str | None = None) -> str:
    """Generate PPT from all existing Phase 3 outputs.

    Scans outputs/ for v3_final.svg + metadata files and builds slides.

    Args:
        output_path: Path for the .pptx file. Defaults to outputs/presentation.pptx.

    Returns:
        Path to the generated PPTX file.
    """
    if output_path is None:
        output_path = str(OUTPUTS_DIR / "presentation.pptx")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_title_slide(prs)

    # Scan outputs directory for samples
    for sample_dir in sorted(OUTPUTS_DIR.iterdir()):
        if not sample_dir.is_dir() or sample_dir.name.startswith("_"):
            continue

        # Find final SVG
        svg_files = sorted(sample_dir.glob("*v3_final.svg"))
        if not svg_files:
            svg_files = sorted(sample_dir.glob("*.svg"))
        if not svg_files:
            continue

        svg_path = svg_files[-1]
        svg_code = svg_path.read_text(encoding="utf-8")

        # Find latest metadata
        meta_files = sorted(sample_dir.glob("metadata_*.json"))
        metadata = {}
        if meta_files:
            with open(meta_files[-1], encoding="utf-8") as f:
                metadata = json.load(f)

        # Find prompt from content IR
        content_ir_path = sample_dir / "01_content_ir.json"
        if not content_ir_path.exists():
            content_ir_path = sample_dir / "content_ir.json"
        prompt = sample_dir.name
        if content_ir_path.exists():
            with open(content_ir_path, encoding="utf-8") as f:
                ir = json.load(f)
            cs = ir.get("content_summary", {})
            prompt = cs.get("title", sample_dir.name)

        _add_sample_slide(prs, sample_dir.name, prompt, metadata, svg_code)
        logger.info(f"Added slide: {sample_dir.name}")

    # ── Summary slide ─────────────────────────────────────
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    summary_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(11), Inches(4)
    )
    tf = summary_box.text_frame
    tf.word_wrap = True

    # Collect all scores
    scores = []
    for sample_dir in sorted(OUTPUTS_DIR.iterdir()):
        if not sample_dir.is_dir() or sample_dir.name.startswith("_"):
            continue
        meta_files = sorted(sample_dir.glob("metadata_*.json"))
        if meta_files:
            with open(meta_files[-1], encoding="utf-8") as f:
                meta = json.load(f)
            ir = meta.get("content_ir_summary", {})
            scores.append({
                "name": sample_dir.name,
                "score": meta.get("final_score", "?"),
                "passed": meta.get("passed", False),
                "intent": ir.get("intent", "?"),
            })

    p = tf.paragraphs[0]
    p.text = "Generation Summary"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    for s in scores:
        p2 = tf.add_paragraph()
        status = "PASS" if s["passed"] else "FAIL"
        p2.text = f"{s['name']}: {s['score']}/10 [{status}] — {s['intent']}"
        p2.font.size = Pt(16)
        p2.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    logger.info(f"PPT saved to {output_path}")
    return output_path


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    path = export_ppt()
    print(f"PPT exported: {path}")
